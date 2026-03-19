from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent.parent
DOCS_DIR = BASE_DIR / "docs"
SCREENSHOTS_DIR = DOCS_DIR / "screenshots"
OUTPUT_PATH = DOCS_DIR / "presentation_soutenance_boutique.pptx"

BG = RGBColor(247, 249, 252)
NAVY = RGBColor(17, 24, 39)
BLUE = RGBColor(37, 99, 235)
MUTED = RGBColor(75, 85, 99)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(226, 232, 240)


def add_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(9.8), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.58), Inches(11.0), Inches(0.25))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(209, 213, 219)


def add_bullets(slide, items, left, top, width, height, font_size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(7)
        first = False
    return box


def add_paragraph(slide, text, left, top, width, height, size=16, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return box


def add_caption(slide, text, left, top, width):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = MUTED


def add_image_card(slide, image_name, left, top, width, height, caption):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SOFT

    image_path = SCREENSHOTS_DIR / image_name
    slide.shapes.add_picture(
        str(image_path),
        left + Inches(0.12),
        top + Inches(0.12),
        width - Inches(0.24),
        height - Inches(0.45),
    )
    add_caption(slide, caption, left, top + height - Inches(0.28), width)


def style_table(table):
    for row in range(len(table.rows)):
        for col in range(len(table.columns)):
            cell = table.cell(row, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if row == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12 if row == 0 else 13)
                    run.font.bold = row == 0
                    run.font.color.rgb = WHITE if row == 0 else NAVY


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(0.85), Inches(1.6), Inches(0.12)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()
    add_paragraph(slide, "Projet Django\nBoutique en ligne", Inches(0.65), Inches(1.15), Inches(8.8), Inches(1.6), size=28, color=WHITE)
    add_paragraph(
        slide,
        "Cours : Management des donnees et innovation\nEquipe : Zeinabou, Assia, Awa, Alexis, Chloe",
        Inches(0.68),
        Inches(3.0),
        Inches(8.2),
        Inches(1.0),
        size=18,
        color=RGBColor(226, 232, 240),
    )
    add_paragraph(
        slide,
        "Soutenance : objectif, fonctionnement, code, demo, difficultes",
        Inches(0.68),
        Inches(4.4),
        Inches(7.8),
        Inches(0.6),
        size=14,
        color=RGBColor(209, 213, 219),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Plan de la presentation")
    add_bullets(
        slide,
        [
            "1. Objectif du projet",
            "2. Fonctionnement du site",
            "3. Structure du code",
            "4. Repartition du travail",
            "5. Choix de design",
            "6. Demonstration",
            "7. Difficultes rencontrees",
            "8. Conclusion",
        ],
        Inches(1.0),
        Inches(1.5),
        Inches(5.2),
        Inches(4.8),
        font_size=22,
    )
    add_paragraph(
        slide,
        "La presentation doit surtout montrer que vous maitrisez l'objectif, la logique du site et votre repartition dans le groupe.",
        Inches(6.6),
        Inches(2.0),
        Inches(5.4),
        Inches(1.8),
        size=18,
        color=MUTED,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Objectif et fonctionnement general")
    add_bullets(
        slide,
        [
            "Developper un site de vente en ligne avec Django et SQLite.",
            "Distinguer deux types d'utilisateurs : client et administrateur.",
            "Permettre la consultation du catalogue, la connexion et l'achat.",
            "Permettre la gestion des articles et la mise a jour du stock.",
        ],
        Inches(0.9),
        Inches(1.5),
        Inches(11.4),
        Inches(4.8),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Ce que le site permet de faire")
    add_bullets(
        slide,
        [
            "Client : voir les articles, les images, le stock et le prix.",
            "Client : se connecter, s'inscrire et acheter une quantite disponible.",
            "Administrateur : ajouter, modifier et supprimer des articles.",
            "Le stock diminue automatiquement apres un achat valide.",
        ],
        Inches(0.9),
        Inches(1.5),
        Inches(5.6),
        Inches(4.8),
    )
    add_bullets(
        slide,
        [
            "Les categories facilitent la navigation.",
            "Le catalogue s'appuie sur les donnees de la base SQLite.",
            "Les images sont chargees depuis static/img.",
            "Les messages de succes et d'erreur guident l'utilisateur.",
        ],
        Inches(6.5),
        Inches(1.5),
        Inches(5.7),
        Inches(4.8),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Structure du code")
    add_bullets(
        slide,
        [
            "boutique_vetements/ : configuration generale Django",
            "shop/ : application principale du projet",
            "templates/ : pages HTML du site",
            "static/ : images, CSS et JavaScript",
            "docs/ : captures, PDF et presentation",
            "scripts/ : automatisation de captures et presentation",
        ],
        Inches(0.9),
        Inches(1.45),
        Inches(5.8),
        Inches(4.9),
        font_size=18,
    )
    add_bullets(
        slide,
        [
            "shop/models.py : modele Article",
            "shop/views.py : logique metier",
            "shop/forms.py : formulaires",
            "shop/urls.py : routes du site",
            "shop/tests.py : tests automatiques",
            "init_data.py : donnees de demonstration",
        ],
        Inches(6.65),
        Inches(1.45),
        Inches(5.6),
        Inches(4.9),
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Quel code fait quoi ?")
    add_bullets(
        slide,
        [
            "models.py : definit la table Article avec designation, representation, stock, prix et categorie.",
            "views.py : gere les pages, l'achat, les redirections et les messages.",
            "forms.py : controle les saisies utilisateur et admin.",
            "templates/ : affichent les pages visibles par le client et l'admin.",
            "admin.py : personnalise l'administration Django.",
            "tests.py : verifie les parcours importants.",
        ],
        Inches(0.9),
        Inches(1.5),
        Inches(11.3),
        Inches(4.9),
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Repartition du travail")
    rows = [
        ("Membre", "Contribution principale"),
        ("Zeinabou", "Structure du projet, URLs, integration finale, verification."),
        ("Assia", "Modele Article, base de donnees, categories, administration."),
        ("Awa", "Logique d'achat, stock, vues et formulaires."),
        ("Alexis", "Templates HTML, CSS, JavaScript, interface visuelle."),
        ("Chloe", "Donnees de demo, tests, captures et documentation."),
    ]
    table = slide.shapes.add_table(6, 2, Inches(0.85), Inches(1.55), Inches(11.7), Inches(4.7)).table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(9.5)
    for row_index, values in enumerate(rows):
        table.cell(row_index, 0).text = values[0]
        table.cell(row_index, 1).text = values[1]
    style_table(table)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Choix de design")
    add_bullets(
        slide,
        [
            "Interface claire et moderne inspiree d'un site e-commerce de vetements.",
            "Palette sobre pour mettre en avant les produits.",
            "Grille de cartes pour lire vite l'image, le prix et le stock.",
            "Navigation simple entre accueil, catalogue, login et gestion.",
        ],
        Inches(0.9),
        Inches(1.5),
        Inches(5.8),
        Inches(4.8),
    )
    add_bullets(
        slide,
        [
            "Nous n'avons pas choisi un design trop charge.",
            "Le but etait d'etre lisible pour la demonstration.",
            "La priorite est restee la fonctionnalite du projet.",
            "Le design soutient la demo sans compliquer l'utilisation.",
        ],
        Inches(6.6),
        Inches(1.5),
        Inches(5.5),
        Inches(4.8),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Demonstration prevue")
    add_bullets(
        slide,
        [
            "Ouvrir l'accueil puis le catalogue.",
            "Montrer les categories homme, femme et enfant.",
            "Se connecter en client avec client1 / client123.",
            "Acheter un article et montrer le message de confirmation.",
            "Verifier que le stock a baisse.",
            "Se connecter en admin avec admin / admin123.",
            "Montrer l'ajout, la modification ou la suppression d'un article.",
        ],
        Inches(0.9),
        Inches(1.45),
        Inches(11.4),
        Inches(5.0),
        font_size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Captures de demonstration client")
    add_image_card(slide, "01-login-page.png", Inches(0.55), Inches(1.45), Inches(3.8), Inches(5.05), "Connexion")
    add_image_card(slide, "02-catalogue-client.png", Inches(4.55), Inches(1.45), Inches(4.2), Inches(5.05), "Catalogue")
    add_image_card(slide, "03-achat-page.png", Inches(8.95), Inches(1.45), Inches(3.8), Inches(5.05), "Achat")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Verification du stock et gestion admin")
    add_image_card(slide, "04-stock-apres-achat.png", Inches(0.75), Inches(1.5), Inches(5.9), Inches(5.0), "Stock apres achat")
    add_image_card(slide, "05-gestion-admin.png", Inches(6.75), Inches(1.5), Inches(5.9), Inches(5.0), "Gestion administrateur")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Difficultes rencontrees et solutions")
    add_bullets(
        slide,
        [
            "Integration de travaux venant de plusieurs membres.",
            "Certains fichiers externes n'etaient pas dans GitHub au debut.",
            "Nettoyage des noms d'images et correspondance avec les categories.",
            "Adaptation de maquettes HTML/CSS dans une vraie structure Django.",
            "Verification du stock pour eviter les erreurs d'achat.",
        ],
        Inches(0.9),
        Inches(1.5),
        Inches(5.8),
        Inches(4.8),
        font_size=17,
    )
    add_bullets(
        slide,
        [
            "Solution : integration finale et verification globale du projet.",
            "Solution : tests automatiques et parcours verifies manuellement.",
            "Solution : rangement des assets et mise a jour des chemins.",
            "Solution : README, PDF, captures et PowerPoint prepares.",
        ],
        Inches(6.6),
        Inches(1.5),
        Inches(5.5),
        Inches(4.8),
        font_size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG)
    add_header_band(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "Le projet repond au cahier des charges du cours.",
            "Les deux roles sont bien distingues : client et administrateur.",
            "Le catalogue, l'achat et la gestion admin fonctionnent.",
            "Le stock est mis a jour automatiquement.",
            "Le projet est documente et pret pour la soutenance.",
        ],
        Inches(0.9),
        Inches(1.55),
        Inches(11.3),
        Inches(3.8),
        font_size=19,
    )
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.6))
    p = footer.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Comptes de demo : admin / admin123   |   client1 / client123"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = BLUE

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
